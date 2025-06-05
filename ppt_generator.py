import os
import json
import requests
from typing import List, Dict
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import time

class PPTGenerator:
    def __init__(self, api_key: str, model_name: str):
        self.client = Groq(api_key=api_key)
        self.model_name = model_name
        
    def generate_presentation_content(self, topic: str, content_names: List[str]) -> Dict:
        """Generate detailed presentation content using Groq API"""
        content_list = "\n".join([f"- {name}" for name in content_names])
        
        prompt = f"""Create a comprehensive professional presentation about "{topic}" with 8-10 slides. 
Each slide should have unique, detailed content with specific information, examples, and data.

CONTENT SECTIONS:
{content_list}

For each content section, provide:
1. Detailed explanation (30-50 words)
2. Key features/points
3. Real-world examples/case studies
4. Best practices
5. Current trends/data

Return ONLY valid JSON with this structure:
{{
    "title": "Presentation Title",
    "slides": [
        {{
            "slide_type": "title",
            "title": "Main Title",
            "subtitle": "Subtitle",
            "background_color": "blue"
        }},
        {{
            "slide_type": "section",
            "title": "Section Title",
            "content": [
                "Detailed point 1 (30-50 words with specific information)",
                "Detailed point 2 (with examples and data)",
                "Detailed point 3 (with practical applications)",
                "Detailed point 4 (with best practices)",
                "Detailed point 5 (with current trends)"
            ],
            "background_color": "white"
        }},
        {{
            "slide_type": "summary",
            "title": "Key Takeaways",
            "content": [
                "Comprehensive summary point 1",
                "Actionable recommendation 2",
                "Strategic insight 3"
            ],
            "background_color": "light_blue"
        }}
    ]
}}

Ensure each bullet point is detailed, unique, and contains substantial information."""

        try:
            response = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system", 
                        "content": "You are an expert presentation designer. Provide detailed, unique content for each slide. Return only valid JSON."
                    },
                    {
                        "role": "user", 
                        "content": prompt
                    }
                ],
                model=self.model_name,
                temperature=0.7,
                max_tokens=4000
            )
            
            content = response.choices[0].message.content.strip()
            
            # Clean JSON extraction
            if content.startswith('```json'):
                content = content.split('```json')[1].split('```')[0].strip()
            elif content.startswith('```'):
                content = content.split('```')[1].split('```')[0].strip()
            
            start_brace = content.find('{')
            end_brace = content.rfind('}')
            if start_brace != -1 and end_brace != -1:
                content = content[start_brace:end_brace+1]
            
            parsed_content = json.loads(content)
            
            if 'slides' not in parsed_content or not parsed_content['slides']:
                raise ValueError("Invalid JSON structure")
            
            return parsed_content
            
        except Exception as e:
            print(f"Error generating content: {e}")
            return self._create_fallback_content(topic, content_names)
    
    def _create_fallback_content(self, topic: str, content_names: List[str]) -> Dict:
        """Create detailed fallback content when API fails"""
        slides = [{
            "slide_type": "title",
            "title": topic,
            "subtitle": "Comprehensive Professional Presentation",
            "background_color": "blue"
        }]
        
        colors = ["white", "light_blue", "white", "gradient"]
        
        for i, content_name in enumerate(content_names):
            slides.append({
                "slide_type": "section",
                "title": content_name,
                "content": [
                    f"Detailed analysis of {content_name}: {topic} involves multiple components including strategic planning and execution frameworks",
                    f"Key features: Core elements include {content_name.lower()} methodology, implementation techniques, and performance metrics",
                    f"Practical example: Case study from Fortune 500 company showing 35% improvement using these methods",
                    f"Best practices: Industry-standard approaches for {content_name.lower()} with measurable results",
                    f"Current trends: 2024 data shows 42% adoption rate in top organizations with positive ROI"
                ],
                "background_color": colors[i % len(colors)]
            })
        
        slides.append({
            "slide_type": "summary",
            "title": "Conclusion & Recommendations",
            "content": [
                "Comprehensive review of all key concepts with supporting data",
                "Actionable 6-month implementation plan with milestones",
                "Strategic roadmap for long-term success and scalability",
                "Resource allocation and team requirements",
                "Q&A and next steps for immediate action"
            ],
            "background_color": "blue"
        })
        
        return {"title": topic, "slides": slides}
    
    def download_image(self, slide_index: int) -> str:
        """Download relevant image based on slide content"""
        try:
            url = f"https://source.unsplash.com/600x400/?{slide_index},business,tech,meeting"
            response = requests.get(url, timeout=10)
            
            if response.status_code == 200:
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
                temp_file.write(response.content)
                temp_file.close()
                return temp_file.name
        except:
            pass
        return None
    
    def get_background_color(self, color_name: str) -> RGBColor:
        """Get professional RGB color"""
        colors = {
            'blue': RGBColor(0, 84, 147),  # Professional blue
            'light_blue': RGBColor(220, 240, 255),
            'white': RGBColor(255, 255, 255),
            'dark': RGBColor(34, 40, 49),
            'gradient': RGBColor(240, 248, 255)
        }
        return colors.get(color_name, colors['white'])
    
    def create_slide(self, presentation: Presentation, slide_data: Dict, slide_index: int):
        """Create a professional slide"""
        slide_type = slide_data.get('slide_type', 'section')
        
        if slide_type == 'title':
            self._create_title_slide(presentation, slide_data)
        elif slide_type == 'summary':
            self._create_summary_slide(presentation, slide_data)
        else:
            self._create_content_slide(presentation, slide_data, slide_index)
    
    def _create_title_slide(self, presentation: Presentation, slide_data: Dict):
        """Create professional title slide"""
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.get_background_color(slide_data.get('background_color', 'blue'))
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get('subtitle', '')
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, presentation: Presentation, slide_data: Dict, slide_index: int):
        """Create detailed content slide"""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.get_background_color(slide_data.get('background_color', 'white'))
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # Content with proper spacing
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            
            # Adjust content area
            content_placeholder.left = Inches(0.5)
            content_placeholder.top = Inches(1.8)
            content_placeholder.width = Inches(6)
            content_placeholder.height = Inches(4.5)
            
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            for i, bullet_point in enumerate(slide_data.get('content', [])):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.space_after = Pt(12)
                p.space_before = Pt(4)
        
        # Add relevant image
        image_path = self.download_image(slide_index)
        if image_path:
            try:
                left = Inches(7)
                top = Inches(2)
                width = Inches(2.5)
                height = Inches(2)
                slide.shapes.add_picture(image_path, left, top, width, height)
                os.unlink(image_path)
            except Exception as e:
                print(f"Image error: {e}")
    
    def _create_summary_slide(self, presentation: Presentation, slide_data: Dict):
        """Create summary/conclusion slide"""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.get_background_color(slide_data.get('background_color', 'light_blue'))
        
        # Title
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # Content with emphasis
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            
            content_placeholder.left = Inches(1)
            content_placeholder.top = Inches(1.8)
            content_placeholder.width = Inches(8)
            content_placeholder.height = Inches(4.5)
            
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, bullet_point in enumerate(slide_data.get('content', [])):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0
                p.font.size = Pt(20)
                p.font.bold = (i == 0)
                p.font.color.rgb = RGBColor(0, 84, 147) if i == 0 else RGBColor(51, 51, 51)
                p.space_after = Pt(14)
    
    def create_powerpoint(self, presentation_data: Dict, output_path: str):
        """Create professional PowerPoint presentation"""
        presentation = Presentation()
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(7.5)
        
        # Create slides
        for i, slide_data in enumerate(presentation_data['slides']):
            try:
                self.create_slide(presentation, slide_data, i)
            except Exception as e:
                print(f"Error creating slide {i+1}: {e}")
        
        presentation.save(output_path)
    
    def create_presentation(self, topic: str, content_names: List[str], output_file: str = None) -> str:
        """Create complete professional PowerPoint presentation"""
        presentation_data = self.generate_presentation_content(topic, content_names)
        
        if not output_file:
            timestamp = int(time.time())
            output_file = f"professional_{topic.replace(' ', '_')}_{timestamp}.pptx"
        
        self.create_powerpoint(presentation_data, output_file)
        return output_file